"""
document_processor.py
─────────────────────
Responsible for:
  1. Accepting uploaded files (PDF, DOCX, TXT, PPTX, XLSX)
  2. Extracting raw text
  3. Semantic chunking via LangChain's RecursiveCharacterTextSplitter

Chunking strategy (RecursiveCharacterTextSplitter)
──────────────────────────────────────────────────
Splits on: paragraph breaks → sentences → words → characters (in order).
This keeps semantically related sentences together in the same chunk,
which significantly improves retrieval quality vs. naive word-count splits.
"""

import re
from pathlib import Path
from typing import List, Dict, Any, Optional

import pypdf
import docx
from pptx import Presentation
import openpyxl
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import CHUNK_SIZE, CHUNK_OVERLAP


# ── Helpers ───────────────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    """
    Normalize whitespace and strip only real control characters.

    Previously used r"[^\x20-\x7E\n]" which destroyed all non-ASCII text
    (accented letters, Arabic, Chinese, smart quotes, emojis, etc.).

    Now strips only C0/C1 control characters (invisible, non-printable)
    while preserving all valid Unicode — including every human language.
    """
    # Collapse multiple spaces / tabs into a single space
    text = re.sub(r"[^\S\n]+", " ", text)
    # Strip only genuine control chars: 0x00-0x08, 0x0B-0x0C, 0x0E-0x1F, 0x7F
    # (keeps \t=0x09, \n=0x0A, \r=0x0D which are meaningful whitespace)
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", text)
    return text.strip()


# ── Extractors ────────────────────────────────────────────────────────────────

def extract_text_from_pdf(path: str) -> str:
    text_parts = []
    with open(path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return _clean_text("\n".join(text_parts))


def extract_text_from_docx(path: str) -> str:
    doc = docx.Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return _clean_text("\n".join(paragraphs))


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return _clean_text(f.read())


def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
    return _clean_text("\n".join(text_parts))


def extract_text_from_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    try:
        text_parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        return _clean_text("\n".join(text_parts))
    finally:
        wb.close()   # always release the file handle


EXTRACTOR_MAP = {
    ".pdf":  extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".txt":  extract_text_from_txt,
    ".pptx": extract_text_from_pptx,
    ".xlsx": extract_text_from_xlsx,
}


def extract_text(file_path: str) -> str:
    """Route to the correct extractor based on file extension."""
    ext = Path(file_path).suffix.lower()
    extractor = EXTRACTOR_MAP.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {list(EXTRACTOR_MAP.keys())}")
    return extractor(file_path)


# ── Chunker ───────────────────────────────────────────────────────────────────

def chunk_text(
    text: str,
    source: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> List[Dict[str, Any]]:
    """
    Split text into semantically coherent, overlapping chunks.

    Uses LangChain's RecursiveCharacterTextSplitter which tries to split
    on paragraph breaks first, then sentences, then words — ensuring chunks
    never slice through the middle of a meaningful sentence.

    Args:
        text:       The full extracted document text.
        source:     Display name embedded in each chunk for source citation.
        chunk_size: Maximum characters per chunk (default from config).
        overlap:    Characters of overlap between consecutive chunks.

    Returns:
        List of dicts: { "text": str, "source": str, "chunk_id": int }
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        # Split order: double-newline (paragraph) → sentence → word → char
        separators=["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""],
        length_function=len,
        is_separator_regex=False,
    )
    raw_chunks = splitter.split_text(text)
    return [
        {"text": chunk, "source": source, "chunk_id": i}
        for i, chunk in enumerate(raw_chunks)
        if chunk.strip()   # drop any blank chunks
    ]


# ── Main entry ────────────────────────────────────────────────────────────────

def process_document(file_path: str, display_name: Optional[str] = None) -> List[Dict[str, Any]]:
    """
    Full pipeline: extract → clean → chunk.

    Args:
        file_path:    Absolute path to the uploaded document on disk.
        display_name: Human-readable filename to embed in chunk metadata
                      (e.g. the original upload name, not the UUID-prefixed
                      stored filename).  Defaults to the basename of file_path.

    Returns:
        List of chunk dicts ready to be embedded.
    """
    filename = display_name or Path(file_path).name
    raw_text = extract_text(file_path)

    if not raw_text:
        raise ValueError(f"Could not extract any text from '{filename}'. The file may be empty or image-based.")

    chunks = chunk_text(raw_text, source=filename)
    return chunks
